# -*- coding: utf-8 -*-
"""
Combined Streamlit Interface and FastAPI Backend for CV Processing
Created on: March 10, 2025
Author: Grok 3 (assisted by xAI)
Description: This script combines a Streamlit web interface and FastAPI backend into a single application,
allowing both interactive CV processing and API-based interactions.
"""
import threading
import time
import streamlit as st
import pdfplumber
import json
import os
import shutil
import openai
from docx import Document
import re
import pandas as pd
from fastapi import FastAPI, HTTPException, Header, File, UploadFile, Form
from fastapi.responses import JSONResponse
from typing import Optional
from openai import OpenAI
import nest_asyncio
import uvicorn
import asyncio
from threading import Thread
from pptx import Presentation
# ============================================ API =========================================================
app = FastAPI()
@app.get("/")
async def root():
    return {"message": "FastAPI is running!"}
def run_fastapi():
    """Chạy FastAPI server trên cổng 8001"""
    uvicorn.run(app, host="0.0.0.0", port=8001, log_level="info")
# Thư mục lưu file sau khi xử lý

TEMP_DIR = "/tmp/"
os.makedirs(TEMP_DIR, exist_ok=True)  # Đảm bảo thư mục tồn tại

# ✅ Hàm làm sạch văn bản
def clean_text(text: str) -> str:
    return re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', text).strip()
def extract_text_from_pptx(file_path: str) -> str:
    """
    Trích xuất toàn bộ văn bản từ file PPTX, gồm text trong slide, shape và table.
    """
    text_content = []
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                # Nếu shape là textbox có text_frame
                if shape.has_text_frame:
                    text_content.append(clean_text(shape.text))

                # Nếu shape chứa table
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = [clean_text(cell.text) for cell in row.cells]
                        text_content.append(" | ".join(row_text))

    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"Lỗi xử lý PPTX: {str(e)}"
        )

    return "\n".join(text_content)
# ✅ Hàm trích xuất văn bản từ file
def extract_text_from_file_api(file_path: str) -> str:
    ext = file_path.split('.')[-1].lower()
    try:
        if ext == "pdf":
            with pdfplumber.open(file_path) as pdf:
                return "\n".join(clean_text(page.extract_text() or "") for page in pdf.pages)
        elif ext == "docx":
            doc = Document(file_path)
            return "\n".join(clean_text(p.text) for p in doc.paragraphs)
        elif ext == "txt":
            with open(file_path, "r", encoding="utf-8") as f:
                return clean_text(f.read())
        elif ext == "xlsx":
            df = pd.read_excel(file_path, dtype=str)
            return clean_text(df.to_string(index=False))
        elif ext == "pptx":
            return extract_text_from_pptx(file_path)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Lỗi xử lý {ext.upper()}: {e}")
    
    raise HTTPException(status_code=400, detail="Định dạng không hỗ trợ")

# ✅ API kiểm tra key
@app.get("/check-key/")
async def check_key(api_key: str = Header(...)):
    """API kiểm tra key có hợp lệ không."""
    if api_key == "valid_key":  # Thay bằng hệ thống kiểm tra thực tế
        return {"valid": True}
    return {"valid": False}

# ✅ API upload file & trích xuất văn bản
@app.post("/upload-file/")
async def upload_file_api(user_key: str = Header(...), file: UploadFile = File(...)):
    """Upload file, trích xuất nội dung và lưu vào thư mục /tmp/."""
    timestamp = int(time.time())  # Định danh duy nhất
    file_ext = file.filename.split('.')[-1]
    unique_filename = f"{user_key}_{timestamp}.{file_ext}"  # Tạo tên file duy nhất
    temp_file_path = os.path.join(TEMP_DIR, unique_filename)

    try:
        # Lưu file gốc
        with open(temp_file_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)

        # Trích xuất văn bản
        extracted_text = extract_text_from_file_api(temp_file_path)

        # Lưu văn bản đã trích xuất vào temp (file .txt)
        extracted_filename = f"{user_key}_{timestamp}.txt"
        extracted_text_path = os.path.join(TEMP_DIR, extracted_filename)
        with open(extracted_text_path, "w", encoding="utf-8") as f:
            f.write(extracted_text)

        return {"message": "File uploaded and processed successfully", "filename": extracted_filename}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Lỗi khi xử lý file: {str(e)}")

# ✅ API lấy văn bản đã trích xuất
@app.get("/get-extract-text/")
async def get_extracted_text_api(user_key: str = Header(...)):
    """Trả về nội dung file đã trích xuất gần nhất của user."""
    files = sorted(
        [f for f in os.listdir(TEMP_DIR) if f.startswith(user_key) and f.endswith(".txt")],
        reverse=True
    )
    
    if not files:
        raise HTTPException(status_code=400, detail="Không tìm thấy file đã trích xuất trước đó.")

    latest_file = files[0]  # Lấy file gần nhất
    with open(os.path.join(TEMP_DIR, latest_file), "r", encoding="utf-8") as f:
        extracted_text = f.read()

    return {"filename": latest_file, "extracted_text": extracted_text}
def clean_json_response(response_text): # ✅ Làm sạch phản hồi GPT - co the dung chung
    """Làm sạch phản hồi GPT để loại bỏ các ký tự không mong muốn trước khi phân tích JSON."""
    response_text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', response_text)
    response_text = response_text.strip()
    if response_text.startswith("```json"):
        response_text = response_text[7:]
    if response_text.endswith("```"):
        response_text = response_text[:-3]
    return response_text.strip()
# ✅ API gửi văn bản đã trích xuất đến OpenAI API và nhận JSON
@app.get("/get-json/")
async def get_json_from_text_api(user_key: str = Header(...), openai_api_key: Optional[str] = Header(None)):
    """Dùng OpenAI để trích xuất JSON từ nội dung đã trích xuất."""
    if not openai_api_key:
        raise HTTPException(status_code=400, detail="Thiếu OpenAI API Key.")

    files = sorted(
        [f for f in os.listdir(TEMP_DIR) if f.startswith(user_key) and f.endswith(".txt")],
        reverse=True
    )

    if not files:
        raise HTTPException(status_code=400, detail="Không tìm thấy file đã trích xuất trước đó.")

    latest_file = files[0]
    with open(os.path.join(TEMP_DIR, latest_file), "r", encoding="utf-8") as f:
        extracted_text = f.read()

    try:
        client = OpenAI(api_key=openai_api_key)
        prompt = f"""
    Trích xuất thông tin từ văn bản CV và trả về JSON hợp lệ:
    {{
        "Name": "",
        "Email": "",
        "Phone": "",
        "Skills": [],  
        "Experience": [],  
        "Education": [],  
        "Certifications": [],  
        "Languages": [],  
        "Strengths": [],  
        "Weaknesses": [],  
        "Additional information": []
    }}
    For the **Languages** field, include:
    - The candidate's native language based on their nationality (e.g., Vietnamese for a candidate from Vietnam).
    - Any foreign language certifications (e.g., TOEIC score) and the corresponding language proficiency level (e.g., English with a proficiency level based on the score).

    For **Strengths and Weaknesses**, analyze the candidate's work experience to identify:
    - **Strengths:** Key skills and attributes demonstrated through their experience.
    - **Weaknesses:** Areas for improvement or challenges faced in their roles.
    Văn bản CV:
    {extracted_text}
    """
        response = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[{"role": "system", "content": "You are an expert in extracting information from CVs (resumes) and images with 10 years of experience in getting the exact information needed to recruit suitable positions for the company."

            "**Context:** I will provide you with resumes of candidates (which can be one or more) or image files containing text."

            "**Your task** is to extract information from the resumes and images I provide (I have taken the text from the resume, and the image will be provided to you below) and return the output as a JSON file."

            "Some of the most important information required for each candidate includes:"
            "- Name"
            "- Email"
            "- Phone number"
            "- Skills"
            "- Experience (including: position, timeline, responsibilities)"
            "- Education (including: degree, institution, timeline, GPA)"
            "- Certifications"
            "- Languages (including proficiency based on nationality and language certifications)"
            "- Strengths (based on the candidate's experience and job description)"
            "- Weaknesses (based on the candidate's experience and job description)"
            "- Additional information (including identification and visa details if provided)"

            "**Task:** Extract the following information from the CV text and return it as JSON."

            "**Output:** JSON file format"

            "***Note:** I can provide you with the text, but in that text will be a synthesis of many resumes of different candidates.*"

            "**REMEMBER:** The output should only be in JSON format."
            },{"role": "user", "content": prompt}]
        )
        json_output = response.choices[0].message.content.strip()
        cleaned_text = clean_json_response(json_output)
        return {"filename": latest_file, "json_output": cleaned_text}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Lỗi khi gọi OpenAI API: {str(e)}")

# ✅ API đặt câu hỏi dựa trên nội dung đã trích xuất
@app.post("/ask-question/")
async def ask_question_api(user_key: str = Header(...), question: str = Form(...), openai_api_key: Optional[str] = Header(None)):
    """Dùng OpenAI để trả lời câu hỏi dựa trên nội dung đã trích xuất."""
    if not openai_api_key:
        raise HTTPException(status_code=400, detail="Thiếu OpenAI API Key.")

    files = sorted(
        [f for f in os.listdir(TEMP_DIR) if f.startswith(user_key) and f.endswith(".txt")],
        reverse=True
    )

    if not files:
        raise HTTPException(status_code=400, detail="Không tìm thấy file đã trích xuất trước đó.")

    latest_file = files[0]
    with open(os.path.join(TEMP_DIR, latest_file), "r", encoding="utf-8") as f:
        extracted_text = f.read()

    try:
        client = OpenAI(api_key=openai_api_key)
        response = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[{"role": "system", "content": f"Play the role of a professional HR, with 10 years of experience in finding potential candidates suitable for the company based on the CV (resume) they send Context: I will provide you with information of each CV (resume) in text form, from which I will ask you some questions related to the CV (resume) of this candidate Task: Please provide the most accurate and closest information to the question I asked, helping me have the most objective view of this candidate so that I can decide whether to hire him or not Tone: solemn, dignified, straightforward, suitable for the office environment, recruitment. Below is the content of the candidate's CV\n{extracted_text}"}, {"role": "user", "content": question}]
        )
        answer = response.choices[0].message.content.strip()
        return {"filename": latest_file, "answer": answer}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Lỗi khi gọi OpenAI API: {str(e)}")

# ✅ API xóa toàn bộ dữ liệu trong /tmp/
@app.delete("/clean/")
async def clean_temp_api():
    """Xóa toàn bộ dữ liệu trong thư mục /tmp/."""
    for file in os.listdir(TEMP_DIR):
        os.remove(os.path.join(TEMP_DIR, file))
    return {"message": "All temporary files deleted successfully"}

# ============================================ API =========================================================

# ========================== STREAMLIT AS STATIC ROUTE ==========================
def run_streamlit():
    """Chạy Streamlit app bên trong FastAPI"""
    os.system("streamlit run interface.py --server.port 8050 --server.headless true --browser.serverAddress 0.0.0.0")

# ========================== CHẠY CẢ FASTAPI VÀ STREAMLIT ==========================
if __name__ == "__main__":
    # Chạy Streamlit trong luồng riêng biệt
    streamlit_thread = threading.Thread(target=run_streamlit, daemon=True)
    streamlit_thread.start()

    # Chạy FastAPI trong luồng chính
    uvicorn.run(app, host="0.0.0.0", port=8080, log_level="info")